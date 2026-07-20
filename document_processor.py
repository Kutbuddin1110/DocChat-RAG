import os
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Routes a file to the right extractor based on its extension.
    Every extractor returns the same shape: [{"page_num", "text"}, ...]
    so chunker.chunk_text() can process any format identically.
    """

    @staticmethod
    def process(file_path: str):
        ext = Path(file_path).suffix.lower()

        handlers = {
            ".pdf": DocumentProcessor._process_pdf,
            ".docx": DocumentProcessor._process_docx,
            ".xlsx": DocumentProcessor._process_excel,
            ".xls": DocumentProcessor._process_excel,
            ".pptx": DocumentProcessor._process_pptx,
            ".txt": DocumentProcessor._process_txt,
        }

        handler = handlers.get(ext)
        if not handler:
            raise ValueError(f"Unsupported file format: {ext}")

        try:
            return handler(file_path)
        except Exception as e:
            logger.error(f"Error processing {file_path}: {e}")
            return None

    @staticmethod
    def _process_pdf(path):
        from chunker import load_pdf_text
        return load_pdf_text(path)

    @staticmethod
    def _process_docx(path):
        from docx import Document
        doc = Document(path)

        sections = []
        current_title = "Section 1"
        current_text = []
        section_count = 1

        def flush():
            text = "\n".join(current_text).strip()
            if text:
                sections.append({"page_num": current_title, "text": text})

        for para in doc.paragraphs:
            is_heading = para.style.name.startswith("Heading") if para.style else False
            if is_heading and current_text:
                flush()
                section_count += 1
                current_title = f"Section {section_count}: {para.text.strip()[:50]}"
                current_text = []
            elif para.text.strip():
                current_text.append(para.text)

        flush()

        # Tables aren't tied to paragraph position by python-docx,
        # so we tag them clearly as a separate block rather than guessing where they go.
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    table_text.append(" | ".join(cells))
        if table_text:
            sections.append({"page_num": "Tables", "text": "\n".join(table_text)})

        return sections if sections else None

    @staticmethod
    def _process_excel(path):
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        pages = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    rows_text.append(" | ".join(str(c) if c is not None else "" for c in row))

            if rows_text:  # skip genuinely empty sheets
                text = f"Sheet: {sheet_name}\n" + "\n".join(rows_text)
                pages.append({"page_num": f"Sheet: {sheet_name}", "text": text})

        return pages if pages else None

    @staticmethod
    def _process_pptx(path):
        from pptx import Presentation
        prs = Presentation(path)
        pages = []

        for i, slide in enumerate(prs.slides, 1):
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                text_parts.append(f"[Speaker notes]: {slide.notes_slide.notes_text_frame.text}")

            if text_parts:
                pages.append({"page_num": f"Slide {i}", "text": "\n".join(text_parts)})

        return pages if pages else None

    @staticmethod
    def _process_txt(path):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return [{"page_num": 1, "text": text}] if text.strip() else None